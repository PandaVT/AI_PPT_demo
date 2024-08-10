# src/ppt/ppt_generator.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
import yaml

class PPTGenerator:
    def __init__(self, config):
        self.config = config
        self.template_path = config['data']['templates_path']
        self.default_background = os.path.join(config['data']['backgroud_path'], config['ppt_generation']['default_backgroud'])
        self.output_path = config['data']['output_path']
        self.slide_title_font_size = Pt(config['ppt_generation']['slide_title_font_size'])
        self.slide_content_font_size = Pt(config['ppt_generation']['slide_content_font_size'])
        self.slide_image_size = config['ppt_generation']['slide_image_size']

    def create_presentation(self, slides_content):
        prs = Presentation()
        for slide_content in slides_content:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Using a blank slide layout
            if slide_content.get('background'):
                self.add_background(slide, slide_content['background'], prs)
            self.add_title(slide, slide_content['title'], slide_content.get('title_position'))
            self.add_content(slide, slide_content['content'], slide_content.get('content_position'))
            if 'image' in slide_content and slide_content['image']:
                self.add_image(slide, slide_content['image'])
        output_file = os.path.join(self.output_path, "generated_presentation.pptx")
        prs.save(output_file)
        return output_file

    def add_background(self, slide, background_path, prs):
        if not os.path.exists(background_path):
            raise FileNotFoundError(f"Background image file not found: {background_path}")
        slide.shapes.add_picture(background_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    def add_title(self, slide, title, position=None):
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
        title_box.text = title
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].font.size = self.slide_title_font_size
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 默认居中偏上显示，如果提供了位置调整，则按位置调整
        if position:
            title_box.left = Inches(position.get('left', 1))
            title_box.top = Inches(position.get('top', 1))

    def add_content(self, slide, content, position=None):
        # 默认在title下面，如果提供了位置调整，则按位置调整
        content_top = position.get('top', 3) if position else 3.5
        content_box = slide.shapes.add_textbox(Inches(1), Inches(content_top), Inches(8), Inches(3))
        content_frame = content_box.text_frame
        content_frame.text = content
        content_frame.paragraphs[0].font.size = self.slide_content_font_size
        content_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        if position:
            content_box.left = Inches(position.get('left', 1))
            content_box.top = Inches(position.get('top', content_top))

    def add_image(self, slide, image_path):
        if not image_path:
            print("No image path provided, skipping image addition.")
            return
        if not os.path.exists(image_path):
            raise FileNotFoundError(f"Image file not found: {image_path}")
        slide.shapes.add_picture(image_path, Inches(1), Inches(2.5), width=Inches(self.slide_image_size[0]/96), height=Inches(self.slide_image_size[1]/96))

# Example usage
if __name__ == "__main__":
    with open('config/config.yaml', 'r') as file:
        config = yaml.safe_load(file)

    ppt_generator = PPTGenerator(config)
    slides_content = [
        {
            "title": "Slide 1 Title",
            "content": "This is the content for slide 1.",
            "background": os.path.join(config['data']['backgroud_path'], "bird.jpg"),
            "title_position": {"left": 1, "top": 0.5},
            "content_position": {"left": 1, "top": 2.5},
            "image": None
        },
        {
            "title": "Slide 2 Title",
            "content": "This is the content for slide 2.",
            "background": None,
            "title_position": {"left": 0, "top": 1.5},
            "content_position": {"left": 0.1, "top": 3},
            "image": os.path.join(config['data']['backgroud_path'], "demo.jpg")
        }
    ]
    output_file = ppt_generator.create_presentation(slides_content)
    print(f"Presentation saved as {output_file}")
