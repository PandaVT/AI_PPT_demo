from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def create_presentation():
    return Presentation()


def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[1]  # 使用标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = title_text
    subtitle.text = subtitle_text


def add_picture_slide(prs, img_path, left_in_inches, top_in_inches, height_in_inches):
    slide_layout = prs.slide_layouts[5]  # 使用图片布局
    slide = prs.slides.add_slide(slide_layout)

    left = Inches(left_in_inches)
    top = Inches(top_in_inches)
    height = Inches(height_in_inches)
    slide.shapes.add_picture(img_path, left, top, height=height)


def add_background_picture_slide_with_text(prs, img_path, title_text, content_text):
    blank_slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(blank_slide_layout)

    # 添加背景图片
    left = top = Inches(0)
    pic = slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)

    # 添加标题
    title_box = slide.shapes.add_textbox(left, top, width=prs.slide_width, height=Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # 添加内容文字
    content_box = slide.shapes.add_textbox(left, Inches(1.5), width=prs.slide_width, height=prs.slide_height - Inches(1.5))
    content_frame = content_box.text_frame
    content_frame.text = content_text
    content_frame.paragraphs[0].font.size = Pt(24)
    content_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)


def save_presentation(prs, file_name):
    prs.save(file_name)


if __name__ == "__main__":
    prs = create_presentation()

    add_title_slide(prs, "Hello, Python-pptx!", "Creating a PowerPoint presentation using Python.")

    img_path = 'demo.jpg'
    add_picture_slide(prs, img_path, 1, 1, 5.5)

    background_img_path = 'demo.jpg'
    add_background_picture_slide_with_text(prs, background_img_path, "Background Title", "This is the content text on a background image slide.")

    save_presentation(prs, 'test_presentation.pptx')
